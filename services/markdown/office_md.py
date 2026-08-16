"""Office (modern OOXML) -> Markdown extractors.

* DOCX -> mammoth converts Word to semantic HTML (Word "Heading N" styles become
  ``<h1>``..``<h6>``, lists/tables preserved) which the shared faithful HTML->MD
  pipeline turns into clean GFM. mammoth's own markdown writer is deprecated;
  HTML + markdownify is its documented recommendation.
* PPTX -> python-pptx: one ``## Slide N`` section per slide (so each slide is a
  chunkable unit), body placeholders (in top-to-bottom reading order) as bullets,
  tables as pipe tables, and speaker notes as a ``### Notes`` subsection.
* XLSX / CSV -> pandas for parsing, then the shared hardened table renderer: one
  ``## SheetName`` + GFM table per sheet (CSV is a single table). Cells are read
  as strings so ids/codes are not reformatted.

Every ZIP-based format (docx/pptx/xlsx) is screened for decompression bombs
before its library reads the archive in-process.
"""

from __future__ import annotations

import io

from .common import MAX_TABLE_ROWS, guard_zip_bomb, join_blocks, rows_to_markdown_table
from .html_md import html_to_markdown


def docx_to_markdown(file_bytes: bytes) -> str:
    """DOCX -> Markdown (mammoth -> semantic HTML -> faithful HTML->MD)."""
    import mammoth

    guard_zip_bomb(file_bytes)
    result = mammoth.convert_to_html(io.BytesIO(file_bytes))
    markdown = html_to_markdown(result.value or "", extract_article=False)
    if not markdown.strip():
        raise ValueError("The document contains no extractable text.")
    return markdown


def pptx_to_markdown(file_bytes: bytes) -> str:
    """PPTX -> Markdown, one section per slide (title / body bullets / tables / notes)."""
    from pptx import Presentation

    guard_zip_bomb(file_bytes)
    prs = Presentation(io.BytesIO(file_bytes))
    blocks: list[str] = []

    for index, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_text = (
            _flatten_soft_breaks(title_shape.text).strip()
            if title_shape is not None and title_shape.text
            else ""
        )
        # python-pptx returns a fresh proxy object per access, so identity
        # (``is``) is unreliable — match the title by its stable shape_id.
        title_id = title_shape.shape_id if title_shape is not None else None
        blocks.append(f"## Slide {index}" + (f": {title_text}" if title_text else ""))

        # Emit body shapes in visual reading order (top-to-bottom, left-to-right)
        # rather than XML/z-order.
        ordered = sorted(
            slide.shapes,
            key=lambda s: ((s.top if s.top is not None else 0), (s.left if s.left is not None else 0)),
        )
        body: list[str] = []
        for shape in ordered:
            if title_id is not None and shape.shape_id == title_id:
                continue
            if shape.has_table:
                rows = [
                    [_flatten_soft_breaks(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                table_md = rows_to_markdown_table(rows)
                if table_md:
                    body.append(table_md)
            elif shape.has_text_frame:
                frame_md = _text_frame_to_markdown(shape.text_frame)
                if frame_md:
                    body.append(frame_md)
        if body:
            blocks.append(join_blocks(body))

        notes = _slide_notes(slide)
        if notes:
            blocks.append("### Notes\n\n" + notes)

    markdown = join_blocks(blocks)
    if not markdown.strip():
        raise ValueError("The presentation contains no extractable text.")
    return markdown


def _flatten_soft_breaks(text: str) -> str:
    """Turn PowerPoint's soft line break into a space.

    python-pptx renders ``<a:br/>`` — the break you get from shift+Enter, which
    PowerPoint decks use constantly for manual line wrapping — as U+000B
    (VERTICAL TAB). It is invisible in a terminal but a raw control character
    in the deliverable, and it survived all the way into the JSONL ``content``
    field. It is a break INSIDE one bullet, so it collapses to a space: a
    newline would split the bullet into an orphan line.
    """
    if not text:
        return text
    return " ".join(text.replace("\v", " ").replace("\f", " ").split())


def _text_frame_to_markdown(text_frame) -> str:
    """Render a text frame's paragraphs as level-indented bullet lines."""
    lines: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = _flatten_soft_breaks(paragraph.text).strip()
        if not text:
            continue
        level = min(getattr(paragraph, "level", 0) or 0, 4)
        lines.append(f"{'  ' * level}- {text}")
    return "\n".join(lines)


def _slide_notes(slide) -> str:
    """Speaker notes as prose. Soft breaks are real line breaks here — notes
    are a paragraph block, not a bullet — so each becomes its own line."""
    try:
        if not slide.has_notes_slide:
            return ""
        notes = slide.notes_slide.notes_text_frame.text
    except Exception:  # noqa: BLE001 — notes are optional metadata
        return ""
    if not notes:
        return ""
    lines = [
        _flatten_soft_breaks(line)
        for line in notes.replace("\v", "\n").replace("\f", "\n").splitlines()
    ]
    return "\n".join(line for line in lines if line).strip()


def xlsx_to_markdown(file_bytes: bytes) -> str:
    """XLSX -> Markdown: one ``## SheetName`` + GFM table per non-empty sheet."""
    import pandas as pd

    guard_zip_bomb(file_bytes)
    # nrows: rows past MAX_TABLE_ROWS are discarded by rows_to_markdown_table
    # anyway, so don't parse them; +1 keeps its "(table truncated)" note firing.
    sheets = pd.read_excel(
        io.BytesIO(file_bytes),
        sheet_name=None,
        engine="openpyxl",
        dtype=str,
        nrows=MAX_TABLE_ROWS + 1,
    )
    blocks: list[str] = []
    for name, frame in sheets.items():
        blocks.append(f"## {name}")
        frame = frame.fillna("")
        if frame.empty:
            blocks.append("_(empty sheet)_")
            continue
        # itertuples streams rows without materialising the whole sheet as one
        # numpy object array (frame.values) — a second full-sheet copy.
        rows = [list(frame.columns), *frame.itertuples(index=False, name=None)]
        table = rows_to_markdown_table(rows)
        blocks.append(table or "_(empty sheet)_")

    markdown = join_blocks(blocks)
    if not markdown.strip():
        raise ValueError("The spreadsheet contains no data.")
    return markdown


def csv_to_markdown(file_bytes: bytes) -> str:
    """CSV -> a single GFM table."""
    import pandas as pd

    from .common import decode_text_bytes

    text = decode_text_bytes(file_bytes)
    try:
        # nrows/itertuples: same bounded-parse + no-full-copy rationale as
        # xlsx_to_markdown above.
        frame = pd.read_csv(
            io.StringIO(text), dtype=str, nrows=MAX_TABLE_ROWS + 1
        ).fillna("")
    except Exception as exc:  # noqa: BLE001 — surface a clean 400
        raise ValueError(f"Could not parse the CSV file: {exc}")
    if frame.empty:
        raise ValueError("The CSV file contains no rows.")
    rows = [list(frame.columns), *frame.itertuples(index=False, name=None)]
    table = rows_to_markdown_table(rows)
    if not table:
        raise ValueError("The CSV file contains no rows.")
    return table
